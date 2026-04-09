# ================================================================
#  APPLICATION STREAMLIT — AUDIT SEO COMPLET
#  Crawler multi-thread + Analyse + PDF + PPTX
#
#  Lancement :
#    pip install streamlit requests beautifulsoup4 pandas
#           matplotlib reportlab python-pptx lxml
#    streamlit run seo_audit_app.py
# ================================================================

import streamlit as st
import requests
from bs4 import BeautifulSoup
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
import threading
import time
import re
import io
import os
import warnings
import xml.etree.ElementTree as ET
import urllib.robotparser
from queue import Queue, Empty
from urllib.parse import urljoin, urlparse
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                 TableStyle, PageBreak, Image as RLImage, HRFlowable)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')

# ─── CONFIG PAGE ───────────────────────────────────────────────
st.set_page_config(
    page_title="SEO Audit Pro",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── CSS PERSONNALISÉ ──────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;600&family=DM+Mono&display=swap');

html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }

.main-header {
    background: linear-gradient(135deg, #1a237e 0%, #283593 50%, #1565c0 100%);
    padding: 2rem 2.5rem;
    border-radius: 16px;
    margin-bottom: 1.5rem;
    color: white;
}
.main-header h1 { margin: 0; font-size: 2rem; font-weight: 600; }
.main-header p  { margin: .4rem 0 0; opacity: .75; font-size: 1rem; }

.metric-card {
    background: white;
    border: 1px solid #e8eaed;
    border-radius: 12px;
    padding: 1.2rem 1.4rem;
    text-align: center;
    box-shadow: 0 1px 3px rgba(0,0,0,.06);
}
.metric-val  { font-size: 2rem; font-weight: 600; margin: 0; }
.metric-lbl  { font-size: .78rem; color: #666; margin: .2rem 0 0; text-transform: uppercase; letter-spacing: .5px; }
.metric-ok   { color: #2e7d32; }
.metric-warn { color: #f57f17; }
.metric-err  { color: #c62828; }
.metric-info { color: #1565c0; }
.metric-gray { color: #455a64; }

.score-bar-wrap { margin: .4rem 0; }
.score-bar-bg   { background: #eceff1; border-radius: 8px; height: 10px; overflow: hidden; }
.score-bar-fill { height: 100%; border-radius: 8px; transition: width .6s ease; }

.issue-row-crit { background: #fff5f5; border-left: 3px solid #c62828; padding: .4rem .8rem; margin: .2rem 0; border-radius: 0 6px 6px 0; font-size: .85rem; }
.issue-row-warn { background: #fffde7; border-left: 3px solid #f9a825; padding: .4rem .8rem; margin: .2rem 0; border-radius: 0 6px 6px 0; font-size: .85rem; }

.status-badge {
    display: inline-block; padding: .15rem .55rem; border-radius: 6px;
    font-size: .75rem; font-weight: 600; font-family: 'DM Mono', monospace;
}
.badge-200 { background: #e8f5e9; color: #2e7d32; }
.badge-3xx { background: #fff8e1; color: #f57f17; }
.badge-4xx { background: #ffebee; color: #c62828; }
.badge-0   { background: #f3e5f5; color: #6a1b9a; }

.section-title {
    font-size: 1.15rem; font-weight: 600; color: #1a237e;
    border-bottom: 2px solid #e3f2fd;
    padding-bottom: .5rem; margin-bottom: 1rem;
}
.crawl-log {
    background: #0d1117; color: #7ee787; font-family: 'DM Mono', monospace;
    font-size: .78rem; padding: 1rem; border-radius: 10px;
    max-height: 280px; overflow-y: auto; line-height: 1.6;
}
.tip-box {
    background: #e3f2fd; border-left: 4px solid #1565c0;
    padding: .7rem 1rem; border-radius: 0 8px 8px 0;
    font-size: .88rem; color: #0d47a1; margin: .5rem 0;
}
</style>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════
#  FONCTIONS MÉTIER (Crawler + Analyse + Export)
# ════════════════════════════════════════════════════════════════

COLORS_MAP = {
    "ok":      "#2e7d32",
    "warn":    "#f57f17",
    "error":   "#c62828",
    "info":    "#1565c0",
    "neutral": "#455a64",
    "bg":      "#f5f5f5",
    "accent":  "#1a237e",
}


def parse_sitemap_urls(base_url: str, headers: dict, max_urls=20000) -> list:
    candidates = [
        f"{base_url.rstrip('/')}/sitemap.xml",
        f"{base_url.rstrip('/')}/sitemap_index.xml",
    ]
    try:
        r = requests.get(f"{base_url.rstrip('/')}/robots.txt", headers=headers, timeout=8)
        for line in r.text.splitlines():
            if line.lower().startswith("sitemap:"):
                u = line.split(":", 1)[1].strip()
                if u not in candidates:
                    candidates.insert(0, u)
    except:
        pass

    to_fetch, found, seen = [], [], set()
    for c in candidates:
        try:
            r = requests.get(c, headers=headers, timeout=12)
            if r.status_code == 200 and ("<url" in r.text or "<sitemap" in r.text):
                to_fetch.append((c, r.text))
                seen.add(c)
                break
        except:
            continue

    def parse_xml(content):
        try:
            root = ET.fromstring(content)
            tag  = root.tag.lower()
            locs = [e.text.strip() for e in root.iter()
                    if e.tag.split("}")[-1] == "loc" and e.text]
            return ("index" if "sitemapindex" in tag else "urls"), locs
        except:
            return "error", []

    while to_fetch and len(found) < max_urls:
        sm_url, content = to_fetch.pop(0)
        kind, data = parse_xml(content)
        if kind == "index":
            for child in data:
                if child not in seen:
                    seen.add(child)
                    try:
                        r2 = requests.get(child, headers=headers, timeout=12)
                        if r2.status_code == 200:
                            to_fetch.append((child, r2.text))
                    except:
                        pass
        elif kind == "urls":
            found.extend(data)

    return list(dict.fromkeys(found))[:max_urls]


def crawl_site_threaded(config: dict,
                         progress_cb=None,
                         log_cb=None,
                         stop_event_ext=None) -> pd.DataFrame:
    """
    Crawler multi-thread avec callbacks pour Streamlit.
    progress_cb(n, total) → met à jour la barre de progression
    log_cb(msg)           → écrit dans le terminal live
    stop_event_ext        → threading.Event pour arrêt externe
    """
    start_url    = config["url"]
    max_pages    = config["max_pages"]
    n_threads    = config.get("threads", 10)
    delay        = config.get("delay", 0.1)
    timeout      = config.get("timeout", 15)

    base_url     = f"{urlparse(start_url).scheme}://{urlparse(start_url).netloc}"
    base_domain  = urlparse(start_url).netloc

    HEADERS = {
        "User-Agent":      "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Accept":          "text/html,application/xhtml+xml,*/*;q=0.8",
        "Accept-Language": "fr-FR,fr;q=0.9,en;q=0.8",
        "Connection":      "keep-alive",
    }

    SKIP_EXT = {'.jpg','.jpeg','.png','.gif','.svg','.webp','.ico',
                '.pdf','.zip','.css','.js','.mp4','.mp3',
                '.woff','.woff2','.ttf','.eot'}

    def is_valid(url):
        try:
            p = urlparse(url)
            if p.netloc != base_domain: return False
            if any(p.path.lower().endswith(e) for e in SKIP_EXT): return False
            skip_q = ('utm_','phpsessid','sort=','filter=','session=')
            if any(s in p.query.lower() for s in skip_q): return False
            return True
        except:
            return False

    # Robots
    rp = None
    if config.get("respect_robots", True):
        rp = urllib.robotparser.RobotFileParser()
        rp.set_url(f"{base_url}/robots.txt")
        try: rp.read()
        except: pass

    def allowed(url):
        if rp is None: return True
        try: return rp.can_fetch("*", url)
        except: return True

    visited_lock = threading.Lock()
    results_lock = threading.Lock()
    visited      = set()
    url_queue    = Queue()
    results      = []
    stop_event   = threading.Event()
    counter      = {"n": 0}
    t_start      = time.time()

    # Alimenter la file
    seeded = 0
    if config.get("use_sitemap", True):
        if log_cb: log_cb("📡 Lecture du sitemap...")
        sm_urls = parse_sitemap_urls(base_url, HEADERS, max_urls=max_pages * 2)
        for u in sm_urls:
            if is_valid(u) and allowed(u):
                url_queue.put(u)
                seeded += 1
        if log_cb: log_cb(f"   ✅ {seeded:,} URLs depuis le sitemap")

    url_queue.put(start_url)

    def extract(url):
        row = {
            "url": url, "status": 0, "load_time": 0,
            "title": "", "title_len": 0,
            "meta_desc": "", "meta_desc_len": 0,
            "h1_count": 0, "h1_text": "",
            "h2_count": 0, "h3_count": 0,
            "canonical": "", "canonical_self": False,
            "noindex": False, "nofollow": False,
            "imgs_total": 0, "imgs_no_alt": 0,
            "internal_links": 0, "external_links": 0,
            "https": url.startswith("https"),
            "content_length": 0,
            "has_schema": False,
            "depth": url.count("/") - 2,
            "error": "",
            "child_urls": [],
        }
        try:
            r = requests.get(url, headers=HEADERS, timeout=timeout, allow_redirects=True)
            row["status"]         = r.status_code
            row["load_time"]      = round(r.elapsed.total_seconds(), 2)
            row["content_length"] = len(r.content)
            if "text/html" not in r.headers.get("content-type",""):
                return row
            if r.status_code != 200:
                return row

            soup = BeautifulSoup(r.text, "lxml")

            t = soup.find("title")
            if t and t.text.strip():
                row["title"]     = t.text.strip()
                row["title_len"] = len(row["title"])

            m = soup.find("meta", attrs={"name": re.compile("^description$", re.I)})
            if m and m.get("content"):
                row["meta_desc"]     = m["content"][:200]
                row["meta_desc_len"] = len(m["content"])

            rob = soup.find("meta", attrs={"name": re.compile("^robots$", re.I)})
            if rob and rob.get("content"):
                c = rob["content"].lower()
                row["noindex"]  = "noindex" in c
                row["nofollow"] = "nofollow" in c

            can = soup.find("link", attrs={"rel": re.compile("^canonical$", re.I)})
            if can and can.get("href"):
                row["canonical"]      = can["href"]
                row["canonical_self"] = can["href"].rstrip("/") == url.rstrip("/")

            h1s = soup.find_all("h1")
            row["h1_count"] = len(h1s)
            row["h2_count"] = len(soup.find_all("h2"))
            row["h3_count"] = len(soup.find_all("h3"))
            if h1s: row["h1_text"] = h1s[0].get_text(strip=True)[:100]

            imgs = soup.find_all("img")
            row["imgs_total"]  = len(imgs)
            row["imgs_no_alt"] = sum(1 for img in imgs if not img.get("alt","").strip())

            child_urls = []
            for a in soup.find_all("a", href=True):
                href = a["href"].strip()
                if not href or href.startswith(("#","mailto:","tel:","javascript:")):
                    continue
                abs_url = urljoin(url, href).split("#")[0]
                clean   = abs_url.split("?")[0]
                if urlparse(clean).netloc == base_domain:
                    row["internal_links"] += 1
                    if is_valid(clean) and allowed(clean):
                        child_urls.append(clean)
                elif abs_url.startswith("http"):
                    row["external_links"] += 1
            row["child_urls"] = list(set(child_urls))

            row["has_schema"] = bool(
                soup.find("script", attrs={"type": "application/ld+json"}) or
                soup.find(attrs={"itemscope": True})
            )
        except requests.exceptions.Timeout:
            row["error"] = "Timeout"
        except Exception as e:
            row["error"] = str(e)[:80]
        return row

    def worker():
        while not stop_event.is_set():
            if stop_event_ext and stop_event_ext.is_set():
                stop_event.set()
                break
            try:
                url = url_queue.get(timeout=3)
            except Empty:
                break

            with visited_lock:
                if url in visited or len(visited) >= max_pages:
                    url_queue.task_done()
                    if len(visited) >= max_pages:
                        stop_event.set()
                    continue
                visited.add(url)
                n = len(visited)

            time.sleep(delay)
            data = extract(url)

            for child in data.get("child_urls", []):
                with visited_lock:
                    if child not in visited and len(visited) < max_pages:
                        url_queue.put(child)

            data.pop("child_urls", None)

            with results_lock:
                results.append(data)
                counter["n"] = len(results)
                n = counter["n"]

                if n % 10 == 0 or n <= 5:
                    elapsed = time.time() - t_start
                    rate    = round(n / elapsed * 60) if elapsed > 0 else 0
                    st_     = "✅" if data["status"]==200 else (
                              "⚠️" if data["status"] in (301,302) else "❌")
                    if log_cb:
                        log_cb(f"[{n:>5}/{max_pages}] {st_} {data['status']} "
                               f"~{rate}p/min — {url[:65]}")
                    if progress_cb:
                        progress_cb(n, max_pages)

            url_queue.task_done()

    threads = []
    for _ in range(n_threads):
        t = threading.Thread(target=worker, daemon=True)
        t.start()
        threads.append(t)

    try:
        url_queue.join()
    except Exception:
        pass
    finally:
        stop_event.set()
        for t in threads:
            t.join(timeout=5)

    elapsed = round(time.time() - t_start, 1)
    if log_cb:
        log_cb(f"\n✅ Terminé en {elapsed}s — {len(results):,} pages crawlées")

    return pd.DataFrame(results)


def analyze_crawl(df: pd.DataFrame, config: dict) -> dict:
    pages_ok = df[df["status"] == 200].copy()

    def score_titles(df):
        if df.empty: return 0, []
        issues, ok = [], 0
        for _, r in df.iterrows():
            t = r.get("title_len", 0)
            if t == 0:
                issues.append({"url": r["url"], "problème": "Title absent", "valeur": "—", "sévérité": "critique"})
            elif t < config["title_min"]:
                issues.append({"url": r["url"], "problème": "Title trop court", "valeur": f"{t} car.", "sévérité": "warning"})
            elif t > config["title_max"]:
                issues.append({"url": r["url"], "problème": "Title trop long", "valeur": f"{t} car.", "sévérité": "warning"})
            else:
                ok += 1
        return int(ok / len(df) * 100) if len(df) else 0, issues

    def score_meta(df):
        if df.empty: return 0, []
        issues, ok = [], 0
        for _, r in df.iterrows():
            m = r.get("meta_desc_len", 0)
            if m == 0:
                issues.append({"url": r["url"], "problème": "Meta desc absente", "valeur": "—", "sévérité": "critique"})
            elif m < config["meta_min"]:
                issues.append({"url": r["url"], "problème": "Meta trop courte", "valeur": f"{m} car.", "sévérité": "warning"})
            elif m > config["meta_max"]:
                issues.append({"url": r["url"], "problème": "Meta trop longue", "valeur": f"{m} car.", "sévérité": "warning"})
            else:
                ok += 1
        return int(ok / len(df) * 100) if len(df) else 0, issues

    def score_h1(df):
        if df.empty: return 0, []
        issues, ok = [], 0
        for _, r in df.iterrows():
            h = r.get("h1_count", 0)
            if h == 0:
                issues.append({"url": r["url"], "problème": "H1 absent", "valeur": "0", "sévérité": "critique"})
            elif h > 1:
                issues.append({"url": r["url"], "problème": f"{h} H1 détectés", "valeur": str(h), "sévérité": "warning"})
            else:
                ok += 1
        return int(ok / len(df) * 100) if len(df) else 0, issues

    def score_images(df):
        p = df[df["imgs_total"] > 0]
        if p.empty: return 100, []
        issues, ok = [], 0
        for _, r in p.iterrows():
            if r.get("imgs_no_alt", 0) > 0:
                issues.append({"url": r["url"], "problème": "Images sans alt", "valeur": f"{r['imgs_no_alt']}", "sévérité": "warning"})
            else:
                ok += 1
        return int(ok / len(p) * 100), issues

    def score_speed(df):
        if df.empty: return 0, []
        issues, ok = [], 0
        for _, r in df.iterrows():
            lt = r.get("load_time", 0)
            if lt >= config["load_time_warn"]:
                issues.append({"url": r["url"], "problème": "Temps élevé", "valeur": f"{lt}s", "sévérité": "warning"})
            else:
                ok += 1
        return int(ok / len(df) * 100) if len(df) else 0, issues

    sc_t, it = score_titles(pages_ok)
    sc_m, im = score_meta(pages_ok)
    sc_h, ih = score_h1(pages_ok)
    sc_i, ii = score_images(pages_ok)
    sc_s, is_= score_speed(pages_ok)

    return {
        "df_pages":       df,
        "df_ok":          pages_ok,
        "total_pages":    len(df),
        "pages_200":      len(pages_ok),
        "pages_errors":   len(df[df["status"].isin([404,500,0])]),
        "pages_redirect": len(df[df["status"].isin([301,302])]),
        "noindex_count":  int(pages_ok["noindex"].sum()) if not pages_ok.empty else 0,
        "no_canonical":   len(pages_ok[pages_ok["canonical"]==""]) if not pages_ok.empty else 0,
        "pages_https":    int(df["https"].sum()),
        "pages_schema":   int(pages_ok["has_schema"].sum()) if not pages_ok.empty else 0,
        "score_global":   int(np.mean([sc_t,sc_m,sc_h,sc_i,sc_s])),
        "scores":         {"Titles":sc_t,"Meta desc":sc_m,"H1":sc_h,"Images alt":sc_i,"Vitesse":sc_s},
        "issues":         {"title":pd.DataFrame(it),"meta":pd.DataFrame(im),
                           "h1":pd.DataFrame(ih),"images":pd.DataFrame(ii),"speed":pd.DataFrame(is_)},
        "status_counts":  df["status"].value_counts().to_dict(),
        "avg_load_time":  round(pages_ok["load_time"].mean(),2) if not pages_ok.empty else 0,
        "avg_title_len":  round(pages_ok["title_len"].mean(),1) if not pages_ok.empty else 0,
        "avg_meta_len":   round(pages_ok["meta_desc_len"].mean(),1) if not pages_ok.empty else 0,
    }


def make_chart_gauge(score):
    color = COLORS_MAP["ok"] if score >= 70 else (COLORS_MAP["warn"] if score >= 40 else COLORS_MAP["error"])
    fig, ax = plt.subplots(figsize=(4,3), subplot_kw=dict(aspect="equal"))
    ax.pie([score,100-score], colors=[color,"#EEEEEE"], startangle=90,
           counterclock=False, wedgeprops=dict(width=0.35))
    ax.text(0,-0.05,f"{score}",ha="center",va="center",fontsize=32,fontweight="bold",color=color)
    ax.text(0,-0.48,"Score /100",ha="center",va="center",fontsize=11,color="#666")
    fig.patch.set_alpha(0); ax.set_facecolor("none")
    fig.tight_layout(); return fig


def make_chart_scores(scores):
    names = list(scores.keys()); vals = list(scores.values())
    clrs  = [COLORS_MAP["ok"] if v>=70 else (COLORS_MAP["warn"] if v>=40 else COLORS_MAP["error"]) for v in vals]
    fig, ax = plt.subplots(figsize=(7,3.5))
    bars = ax.barh(names, vals, color=clrs, height=0.5, edgecolor="none")
    for bar,val in zip(bars,vals):
        ax.text(bar.get_width()+1, bar.get_y()+bar.get_height()/2,
                f"{val}", va="center", fontsize=11, fontweight="bold", color="#333")
    ax.axvline(70,color="#999",linestyle="--",linewidth=1,alpha=0.6,label="Seuil 70")
    ax.set_xlim(0,112); ax.set_xlabel("Score /100"); ax.legend(fontsize=9)
    ax.spines[["top","right"]].set_visible(False)
    ax.grid(axis="x",alpha=0.3); fig.tight_layout(); return fig


def make_chart_http(status_counts):
    labels_m = {200:"200 OK",301:"301",302:"302",404:"404",500:"500",0:"Erreur"}
    clrs_m   = {200:COLORS_MAP["ok"],301:COLORS_MAP["info"],302:"#64B5F6",
                404:COLORS_MAP["error"],500:"#880E4F",0:"#BF360C"}
    vals     = [float(v) for v in status_counts.values()]
    lbls     = [labels_m.get(k,str(k)) for k in status_counts.keys()]
    clrs     = [clrs_m.get(k,"#9E9E9E") for k in status_counts.keys()]
    fig, ax  = plt.subplots(figsize=(4,4))
    ax.pie(vals, labels=lbls, colors=clrs,
           autopct=lambda p: f"{p:.0f}%" if p > 3 else "",
           startangle=90, textprops={"fontsize":10})
    fig.tight_layout(); return fig


def make_chart_titles(df_ok, title_min, title_max):
    tl = df_ok["title_len"].dropna()
    fig, ax = plt.subplots(figsize=(7,3))
    ax.hist(tl, bins=20, color=COLORS_MAP["info"], alpha=0.8, edgecolor="white")
    ax.axvline(title_min,color=COLORS_MAP["warn"],linestyle="--",linewidth=1.5,label=f"Min {title_min}")
    ax.axvline(title_max,color=COLORS_MAP["error"],linestyle="--",linewidth=1.5,label=f"Max {title_max}")
    ax.set_xlabel("Longueur (caractères)"); ax.set_ylabel("Nb pages")
    ax.spines[["top","right"]].set_visible(False); ax.grid(alpha=0.3)
    ax.legend(fontsize=9); fig.tight_layout(); return fig


def make_chart_speed(df_ok, threshold):
    lt = df_ok["load_time"].dropna().sort_values()
    fig, ax = plt.subplots(figsize=(7,3))
    clrs = [COLORS_MAP["ok"] if t < threshold else COLORS_MAP["error"] for t in lt]
    ax.bar(range(len(lt)), lt, color=clrs, width=0.8)
    ax.axhline(threshold, color=COLORS_MAP["error"], linestyle="--",
               linewidth=1.5, label=f"Seuil {threshold}s")
    ax.set_xlabel("Pages (triées)"); ax.set_ylabel("Secondes")
    ax.spines[["top","right"]].set_visible(False); ax.grid(axis="y",alpha=0.3)
    ax.legend(fontsize=9); fig.tight_layout(); return fig


def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0); plt.close(fig); return buf


def generate_pdf_report(audit, config) -> bytes:
    buf    = io.BytesIO()
    doc    = SimpleDocTemplate(buf, pagesize=A4,
                                rightMargin=2*cm, leftMargin=2*cm,
                                topMargin=2.5*cm, bottomMargin=2*cm)
    W, H   = A4
    base   = getSampleStyleSheet()

    S = {
        "cover_h1": ParagraphStyle("ch1", fontSize=26, fontName="Helvetica-Bold",
                     textColor=colors.HexColor("#1a237e"), alignment=TA_CENTER, spaceAfter=6),
        "cover_sub": ParagraphStyle("cs", fontSize=14, fontName="Helvetica",
                     textColor=colors.HexColor("#455a64"), alignment=TA_CENTER, spaceAfter=4),
        "section": ParagraphStyle("sec", fontSize=15, fontName="Helvetica-Bold",
                    textColor=colors.HexColor("#1a237e"), spaceBefore=16, spaceAfter=8),
        "sub": ParagraphStyle("sub", fontSize=12, fontName="Helvetica-Bold",
                textColor=colors.HexColor("#455a64"), spaceBefore=10, spaceAfter=5),
        "body": ParagraphStyle("bd", fontSize=10, fontName="Helvetica",
                 textColor=colors.HexColor("#333"), leading=15, spaceAfter=5),
        "small": ParagraphStyle("sm", fontSize=8.5, fontName="Helvetica",
                  textColor=colors.HexColor("#555"), leading=12),
        "ok":   ParagraphStyle("ok",  fontSize=10, fontName="Helvetica-Bold",
                 textColor=colors.HexColor("#2e7d32")),
        "err":  ParagraphStyle("err", fontSize=10, fontName="Helvetica-Bold",
                 textColor=colors.HexColor("#c62828")),
        "cap":  ParagraphStyle("cap", fontSize=9, fontName="Helvetica-Oblique",
                 textColor=colors.HexColor("#777"), alignment=TA_CENTER),
    }

    def rl_img(fig, w=14*cm, h=6.5*cm):
        b = fig_to_buf(fig); b.seek(0)
        return RLImage(b, width=w, height=h)

    def kpi_row(kpis):
        cw = (W - 4*cm) / len(kpis)
        cells = [[
            Paragraph(f"<b>{v}</b>",
                ParagraphStyle("kv", fontSize=20, fontName="Helvetica-Bold",
                    textColor=colors.HexColor(c), alignment=TA_CENTER)),
            Paragraph(l, ParagraphStyle("kl", fontSize=8.5, fontName="Helvetica",
                textColor=colors.HexColor("#666"), alignment=TA_CENTER))
        ] for l,v,c in kpis]
        val_row = [c[0] for c in cells]
        lbl_row = [c[1] for c in cells]
        t = Table([val_row, lbl_row], colWidths=[cw]*len(kpis))
        t.setStyle(TableStyle([
            ("BACKGROUND",(0,0),(-1,-1),colors.HexColor("#f5f5f5")),
            ("ALIGN",(0,0),(-1,-1),"CENTER"),
            ("TOPPADDING",(0,0),(-1,-1),8),
            ("BOTTOMPADDING",(0,0),(-1,-1),8),
        ])); return t

    def issues_tbl(df_iss):
        if df_iss is None or df_iss.empty:
            return Paragraph("✅ Aucun problème détecté.", S["ok"])
        rows = [[Paragraph(c, S["small"]) for c in ["URL","Problème","Valeur","Sévérité"]]]
        for _, r in df_iss.head(25).iterrows():
            sev_col = "#c62828" if r.get("sévérité")=="critique" else "#f57f17"
            rows.append([
                Paragraph(str(r.get("url",""))[-60:], S["small"]),
                Paragraph(str(r.get("problème","")), S["small"]),
                Paragraph(str(r.get("valeur","")), S["small"]),
                Paragraph(str(r.get("sévérité","")),
                    ParagraphStyle("sv",fontSize=8.5,fontName="Helvetica-Bold",
                        textColor=colors.HexColor(sev_col))),
            ])
        cws = [(W-4*cm)*f for f in [.48,.28,.12,.12]]
        t = Table(rows, colWidths=cws, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#1a237e")),
            ("TEXTCOLOR",(0,0),(-1,0),colors.white),
            ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f5f5f5")]),
            ("GRID",(0,0),(-1,-1),.3,colors.HexColor("#ddd")),
            ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
            ("LEFTPADDING",(0,0),(-1,-1),5),
        ])); return t

    domain = urlparse(config["url"]).netloc
    sc     = audit["score_global"]
    story  = []

    # Cover
    story += [Spacer(1,3*cm),
              Paragraph("RAPPORT D'AUDIT SEO", S["cover_h1"]),
              Spacer(1,.3*cm),
              HRFlowable(width="50%",thickness=2,color=colors.HexColor("#1a237e"),hAlign="CENTER"),
              Spacer(1,.5*cm),
              Paragraph(domain, ParagraphStyle("dom",fontSize=18,fontName="Helvetica-Bold",
                  textColor=colors.HexColor("#1565c0"),alignment=TA_CENTER)),
              Spacer(1,1.2*cm),
              rl_img(make_chart_gauge(sc), w=7*cm, h=4.5*cm),
              Spacer(1,1.2*cm),
              Paragraph(f"Client : {config.get('client_name','—')}", S["cover_sub"]),
              Paragraph(f"Consultant : {config.get('consultant','—')}", S["cover_sub"]),
              Paragraph(f"Date : {datetime.now().strftime('%d/%m/%Y')}", S["cover_sub"]),
              PageBreak()]

    # Résumé
    story += [Paragraph("1. Résumé exécutif", S["section"]),
              HRFlowable(width="100%",thickness=1,color=colors.HexColor("#1a237e")),
              Spacer(1,.3*cm)]
    verdict = ("Excellent" if sc>=80 else "Satisfaisant" if sc>=60 else "À améliorer" if sc>=40 else "Critique")
    story += [Paragraph(f"Audit de <b>{domain}</b> — {audit['total_pages']} pages crawlées. "
                        f"Score global : <b>{sc}/100</b> ({verdict}).", S["body"]),
              Spacer(1,.3*cm),
              kpi_row([
                  ("Pages crawlées",  str(audit["total_pages"]),  "#455a64"),
                  ("200 OK",          str(audit["pages_200"]),    "#2e7d32"),
                  ("Erreurs",         str(audit["pages_errors"]), "#c62828"),
                  ("Redirections",    str(audit["pages_redirect"]),"#f57f17"),
              ]),
              Spacer(1,.4*cm),
              kpi_row([
                  ("Noindex",         str(audit["noindex_count"]), "#f57f17"),
                  ("Sans canonical",  str(audit["no_canonical"]),  "#f57f17"),
                  ("Schema.org",      str(audit["pages_schema"]),  "#1565c0"),
                  ("Tps moyen",       f"{audit['avg_load_time']}s","#455a64"),
              ]),
              Spacer(1,.5*cm),
              Paragraph("Scores par axe", S["sub"]),
              rl_img(make_chart_scores(audit["scores"]), w=14*cm, h=5*cm),
              PageBreak()]

    # Technique
    story += [Paragraph("2. Analyse technique", S["section"]),
              HRFlowable(width="100%",thickness=1,color=colors.HexColor("#1a237e")),
              Spacer(1,.3*cm),
              Paragraph("2.1 Codes HTTP", S["sub"]),
              rl_img(make_chart_http(audit["status_counts"]), w=9*cm, h=6.5*cm),
              Spacer(1,.4*cm),
              Paragraph("2.2 Temps de chargement", S["sub"]),
              rl_img(make_chart_speed(audit["df_ok"], config["load_time_warn"]), w=14*cm, h=5*cm)]

    if not audit["issues"]["speed"].empty:
        story += [Spacer(1,.3*cm), issues_tbl(audit["issues"]["speed"])]
    story.append(PageBreak())

    # On-page
    story += [Paragraph("3. On-Page SEO", S["section"]),
              HRFlowable(width="100%",thickness=1,color=colors.HexColor("#1a237e"))]

    for axe_key, axe_name, chart_fn in [
        ("title",  "Balises Title",         lambda: rl_img(make_chart_titles(audit["df_ok"], config["title_min"], config["title_max"]), w=13*cm, h=4.5*cm)),
        ("meta",   "Meta descriptions",     None),
        ("h1",     "Balises H1",            None),
        ("images", "Attributs Alt images",  None),
    ]:
        sc_axe = {"title":audit["scores"]["Titles"],"meta":audit["scores"]["Meta desc"],
                  "h1":audit["scores"]["H1"],"images":audit["scores"]["Images alt"]}[axe_key]
        story += [Spacer(1,.5*cm),
                  Paragraph(f"3.x {axe_name} — Score : {sc_axe}/100", S["sub"])]
        if chart_fn:
            story.append(chart_fn())
        story.append(issues_tbl(audit["issues"][axe_key]))

    story.append(PageBreak())

    # Plan d'action
    story += [Paragraph("4. Plan d'action prioritaire", S["section"]),
              HRFlowable(width="100%",thickness=1,color=colors.HexColor("#1a237e")),
              Spacer(1,.4*cm)]

    actions_sorted = sorted([
        ("Titles",    audit["scores"]["Titles"],    "Corriger titles manquants/hors limites"),
        ("Meta desc", audit["scores"]["Meta desc"], "Rédiger meta descriptions uniques"),
        ("H1",        audit["scores"]["H1"],        "1 seul H1 par page"),
        ("Images",    audit["scores"]["Images alt"],"Attribut alt sur toutes les images"),
        ("Vitesse",   audit["scores"]["Vitesse"],   "Optimiser temps de chargement"),
    ], key=lambda x: x[1])

    reco = [[Paragraph(c, S["small"]) for c in ["P","Axe","Action","Impact"]]]
    for i,(axe,score,action) in enumerate(actions_sorted,1):
        col = "#c62828" if score<40 else "#f57f17"
        reco.append([
            Paragraph(f"P{i}", ParagraphStyle("pn",fontSize=11,fontName="Helvetica-Bold",
                textColor=colors.HexColor(col),alignment=TA_CENTER)),
            Paragraph(axe, S["small"]),
            Paragraph(action, S["small"]),
            Paragraph("Élevé" if score<60 else "Moyen", S["small"]),
        ])
    t = Table(reco, colWidths=[(W-4*cm)*f for f in [.08,.15,.62,.15]], repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#1a237e")),
        ("TEXTCOLOR",(0,0),(-1,0),colors.white),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f5f5f5")]),
        ("GRID",(0,0),(-1,-1),.3,colors.HexColor("#ddd")),
        ("TOPPADDING",(0,0),(-1,-1),6),("BOTTOMPADDING",(0,0),(-1,-1),6),
        ("LEFTPADDING",(0,0),(-1,-1),6),("ALIGN",(0,0),(0,-1),"CENTER"),
    ]))
    story.append(t)

    doc.build(story)
    return buf.getvalue()


def generate_pptx_report(audit, config) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    domain = urlparse(config["url"]).netloc

    C = {
        "dark":  RGBColor(0x1A,0x23,0x7E),
        "blue":  RGBColor(0x15,0x65,0xC0),
        "ok":    RGBColor(0x2E,0x7D,0x32),
        "warn":  RGBColor(0xF5,0x7F,0x17),
        "err":   RGBColor(0xC6,0x28,0x28),
        "white": RGBColor(0xFF,0xFF,0xFF),
        "light": RGBColor(0xF5,0xF5,0xF5),
        "gray":  RGBColor(0x45,0x5A,0x64),
    }

    def add_rect(slide,l,t,w,h,fill=None):
        shp = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        shp.line.fill.background()
        if fill: shp.fill.solid(); shp.fill.fore_color.rgb=fill
        else:    shp.fill.background()
        return shp

    def add_text(slide,text,l,t,w,h,size=16,bold=False,color=None,align=PP_ALIGN.LEFT):
        txb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        txb.word_wrap=True; tf=txb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=align
        run=p.add_run(); run.text=str(text)
        run.font.size=Pt(size); run.font.bold=bold
        if color: run.font.color.rgb=color
        return txb

    def add_img(slide,fig,l,t,w,h):
        b=fig_to_buf(fig); b.seek(0)
        slide.shapes.add_picture(b,Inches(l),Inches(t),Inches(w),Inches(h))

    def header(slide,title,sub=""):
        add_rect(slide,0,0,13.33,1.1,fill=C["dark"])
        add_text(slide,title,.4,.1,10,.6,size=22,bold=True,color=C["white"])
        if sub:
            add_text(slide,sub,.4,.66,10,.35,size=11,
                     color=RGBColor(0xB0,0xBE,0xC5))
        add_text(slide,domain,10.5,.3,2.6,.5,size=10,
                 color=RGBColor(0x90,0xA4,0xAE),align=PP_ALIGN.RIGHT)

    def sc_col(s):
        return C["ok"] if s>=70 else (C["warn"] if s>=40 else C["err"])

    # Slide 1 — Cover
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl,0,0,13.33,7.5,fill=C["dark"])
    add_rect(sl,0,4.8,13.33,2.7,fill=C["blue"])
    add_text(sl,"AUDIT SEO",1,1.2,11.33,1.4,size=44,bold=True,
             color=C["white"],align=PP_ALIGN.CENTER)
    add_text(sl,domain,1,2.7,11.33,.8,size=24,
             color=RGBColor(0x90,0xCA,0xF9),align=PP_ALIGN.CENTER)
    add_text(sl,f"Score global : {audit['score_global']}/100",
             1,3.6,11.33,.7,size=20,bold=True,
             color=sc_col(audit["score_global"]),align=PP_ALIGN.CENTER)
    add_text(sl,f"Client : {config.get('client_name','—')}   ·   "
               f"{datetime.now().strftime('%d/%m/%Y')}",
             1,5,11.33,.5,size=14,color=C["white"],align=PP_ALIGN.CENTER)

    # Slide 2 — Score global
    sl = prs.slides.add_slide(BLANK)
    header(sl,"Score global SEO","Vue d'ensemble")
    add_img(sl,make_chart_gauge(audit["score_global"]),.5,1.2,5,3.8)
    add_img(sl,make_chart_scores(audit["scores"]),5.8,1.2,7.2,4)

    # Slide 3 — KPIs
    sl = prs.slides.add_slide(BLANK)
    header(sl,"Chiffres clés","Résumé du crawl")
    kpis = [
        (audit["total_pages"],"Pages crawlées",C["gray"]),
        (audit["pages_200"],"200 OK",C["ok"]),
        (audit["pages_errors"],"Erreurs",C["err"]),
        (audit["pages_redirect"],"Redirections",C["warn"]),
        (audit["noindex_count"],"Noindex",C["warn"]),
        (audit["no_canonical"],"Sans canonical",C["warn"]),
        (audit["pages_schema"],"Schema.org",C["blue"]),
        (f"{audit['avg_load_time']}s","Tps charg. moy.",C["gray"]),
    ]
    for i,(val,lbl,col) in enumerate(kpis):
        ci,ri = i%4, i//4
        lx,ty = .5+ci*3.1, 1.35+ri*2.55
        add_rect(sl,lx,ty,2.9,2.2,fill=C["light"])
        add_text(sl,str(val),lx,ty+.15,2.9,1.2,size=32,bold=True,
                 color=col,align=PP_ALIGN.CENTER)
        add_text(sl,lbl,lx,ty+1.45,2.9,.55,size=11,
                 color=C["gray"],align=PP_ALIGN.CENTER)

    # Slide 4 — Codes HTTP
    sl = prs.slides.add_slide(BLANK)
    header(sl,"Codes HTTP","Statuts des pages crawlées")
    add_img(sl,make_chart_http(audit["status_counts"]),.5,1.2,5.5,5.5)
    ty=1.4
    add_rect(sl,6.3,ty-.05,6.5,.45,fill=C["dark"])
    for lx,txt in [(6.5,"Code"),(8,"Nb"),(9.2,"Signification")]:
        add_text(sl,txt,lx,ty,.8+2,.38,size=11,bold=True,color=C["white"])
    ty+=.45
    lbls_m={200:"OK",301:"Redir. perm.",302:"Redir. temp.",404:"Non trouvé",500:"Erreur srv",0:"Réseau err."}
    col_m={200:C["ok"],301:C["blue"],302:C["blue"],404:C["err"],500:C["err"],0:C["err"]}
    for j,(code,count) in enumerate(sorted(audit["status_counts"].items())):
        bg=C["light"] if j%2==0 else C["white"]
        add_rect(sl,6.3,ty,6.5,.42,fill=bg)
        add_text(sl,str(code),6.5,ty+.04,1,.34,size=11,bold=True,color=col_m.get(code,C["gray"]))
        add_text(sl,str(count),8,ty+.04,1,.34,size=11,color=C["gray"])
        add_text(sl,lbls_m.get(code,"-"),9.2,ty+.04,3.4,.34,size=10,color=C["gray"])
        ty+=.42

    # Slides axes on-page
    axes = [
        ("Balises Title",       "title",  audit["scores"]["Titles"]),
        ("Meta descriptions",   "meta",   audit["scores"]["Meta desc"]),
        ("Balises H1",          "h1",     audit["scores"]["H1"]),
        ("Images — Alt",        "images", audit["scores"]["Images alt"]),
    ]
    for axe_name,key,score in axes:
        sl = prs.slides.add_slide(BLANK)
        header(sl,f"On-Page — {axe_name}")
        sc_c = sc_col(score)
        add_rect(sl,.5,1.3,3.3,2.3,fill=C["light"])
        add_text(sl,f"{score}/100",.5,1.4,3.3,1.5,size=44,bold=True,
                 color=sc_c,align=PP_ALIGN.CENTER)
        verdict = "Bon" if score>=70 else ("Moyen" if score>=40 else "Critique")
        add_text(sl,verdict,.5,3,3.3,.5,size=16,color=sc_c,align=PP_ALIGN.CENTER)

        df_iss = audit["issues"].get(key, pd.DataFrame())
        if not df_iss.empty:
            n_crit = len(df_iss[df_iss["sévérité"]=="critique"]) if "sévérité" in df_iss.columns else 0
            add_text(sl,f"{len(df_iss)} problème(s)",4.2,1.4,8.5,.5,size=15,bold=True,color=C["err"])
            add_text(sl,f"  Critiques : {n_crit}",4.2,2,8.5,.4,size=13,color=C["err"])
            add_text(sl,f"  Avertissements : {len(df_iss)-n_crit}",4.2,2.5,8.5,.4,size=13,color=C["warn"])
            ty2=3.2
            add_rect(sl,4.2,ty2-.05,8.5,.42,fill=C["dark"])
            add_text(sl,"URL",4.3,ty2,5.5,.38,size=10,bold=True,color=C["white"])
            add_text(sl,"Problème",9.8,ty2,3,.38,size=10,bold=True,color=C["white"])
            ty2+=.42
            for _,row in df_iss.head(8).iterrows():
                bg=C["light"] if _%2==0 else C["white"]
                add_rect(sl,4.2,ty2,8.5,.38,fill=bg)
                add_text(sl,str(row.get("url",""))[-52:],4.3,ty2+.02,5.4,.34,size=8.5,color=C["gray"])
                add_text(sl,str(row.get("problème","")),9.8,ty2+.02,2.8,.34,size=8.5,color=C["err"])
                ty2+=.38
        else:
            add_text(sl,"✅ Aucun problème",4.2,2.5,8.5,.6,size=18,color=C["ok"],align=PP_ALIGN.CENTER)

    # Slide plan d'action
    sl = prs.slides.add_slide(BLANK)
    header(sl,"Plan d'action prioritaire","Actions classées par impact")
    actions = sorted([
        ("Titles",   audit["scores"]["Titles"],   "Corriger titles manquants/hors limites"),
        ("Meta",     audit["scores"]["Meta desc"], "Rédiger meta descriptions uniques"),
        ("H1",       audit["scores"]["H1"],        "1 seul H1 par page"),
        ("Images",   audit["scores"]["Images alt"],"Attribut alt sur toutes les images"),
        ("Vitesse",  audit["scores"]["Vitesse"],   "Optimiser les temps de chargement"),
    ], key=lambda x:x[1])
    ty=1.4
    add_rect(sl,.4,ty-.05,12.5,.42,fill=C["dark"])
    for lx,t in [(.5,"P"),(1.5,"Axe"),(3.2,"Action"),(10.7,"Impact")]:
        add_text(sl,t,lx,ty,1.5,.38,size=11,bold=True,color=C["white"])
    ty+=.42
    for j,(axe,score,action) in enumerate(actions):
        bg=C["light"] if j%2==0 else C["white"]
        add_rect(sl,.4,ty,12.5,.5,fill=bg)
        col_p=C["err"] if score<40 else C["warn"]
        add_text(sl,f"P{j+1}",.5,ty+.05,1,.4,size=14,bold=True,color=col_p,align=PP_ALIGN.CENTER)
        add_text(sl,axe,1.5,ty+.05,1.5,.4,size=11,color=C["dark"])
        add_text(sl,action,3.2,ty+.05,7.2,.4,size=10.5,color=C["gray"])
        add_text(sl,"Élevé" if score<60 else "Moyen",10.7,ty+.05,1.8,.4,size=11,
                 bold=True,color=C["err"] if score<60 else C["warn"],align=PP_ALIGN.CENTER)
        ty+=.5

    # Slide finale
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl,0,0,13.33,7.5,fill=C["dark"])
    add_text(sl,"Merci",1,2.2,11.33,1.5,size=52,bold=True,color=C["white"],align=PP_ALIGN.CENTER)
    add_text(sl,config.get("consultant",""),1,4,11.33,.7,size=18,
             color=RGBColor(0x90,0xCA,0xF9),align=PP_ALIGN.CENTER)
    add_text(sl,config["url"],1,4.9,11.33,.6,size=13,
             color=RGBColor(0x78,0x90,0x9C),align=PP_ALIGN.CENTER)

    out = io.BytesIO(); prs.save(out); return out.getvalue()


# ════════════════════════════════════════════════════════════════
#  INTERFACE STREAMLIT
# ════════════════════════════════════════════════════════════════

# Header
st.markdown("""
<div class="main-header">
  <h1>🔍 SEO Audit Pro</h1>
  <p>Crawler multi-thread · Analyse on-page · Rapports PDF & Google Slides</p>
</div>
""", unsafe_allow_html=True)


# ─── SIDEBAR — Configuration ────────────────────────────────────
with st.sidebar:
    st.markdown("### ⚙️ Configuration")

    url_input = st.text_input(
        "URL du site à auditer",
        value="https://docs.streamlit.io",
        placeholder="https://exemple.fr",
    )

    with st.expander("🕷️ Paramètres du crawl", expanded=True):
        max_pages = st.number_input("Max pages", min_value=1, max_value=500_000,
                                     value=100, step=50)
        n_threads = st.slider("Threads parallèles", 1, 50, 8)
        delay     = st.slider("Délai entre requêtes (s)", 0.0, 2.0, 0.1, step=0.05)
        col1, col2 = st.columns(2)
        use_sitemap    = col1.checkbox("Sitemap XML", value=True)
        respect_robots = col2.checkbox("Robots.txt", value=True)

    with st.expander("📏 Seuils d'alerte"):
        col1, col2 = st.columns(2)
        title_min = col1.number_input("Title min (car.)", value=30)
        title_max = col2.number_input("Title max (car.)", value=60)
        col3, col4 = st.columns(2)
        meta_min  = col3.number_input("Meta min (car.)",  value=70)
        meta_max  = col4.number_input("Meta max (car.)",  value=160)
        load_warn = st.slider("Seuil temps charg. (s)", 0.5, 10.0, 3.0, step=0.5)

    with st.expander("📄 Infos rapport"):
        client_name = st.text_input("Nom du client", value="Mon Client")
        consultant  = st.text_input("Consultant",     value="Votre Nom")

    st.markdown("---")
    btn_start = st.button("🚀 Lancer l'audit", type="primary", use_container_width=True)
    btn_stop  = st.button("⏹ Arrêter", use_container_width=True)


# ─── ÉTAT SESSION ────────────────────────────────────────────────
if "audit_done"   not in st.session_state: st.session_state.audit_done   = False
if "audit_result" not in st.session_state: st.session_state.audit_result = None
if "df_pages"     not in st.session_state: st.session_state.df_pages     = None
if "stop_event"   not in st.session_state: st.session_state.stop_event   = None
if "config"       not in st.session_state: st.session_state.config       = {}

config = {
    "url":            url_input,
    "max_pages":      int(max_pages),
    "threads":        int(n_threads),
    "delay":          float(delay),
    "timeout":        15,
    "use_sitemap":    use_sitemap,
    "respect_robots": respect_robots,
    "title_min":      int(title_min),
    "title_max":      int(title_max),
    "meta_min":       int(meta_min),
    "meta_max":       int(meta_max),
    "load_time_warn": float(load_warn),
    "client_name":    client_name,
    "consultant":     consultant,
}


# ─── LANCEMENT DU CRAWL ──────────────────────────────────────────
if btn_stop and st.session_state.stop_event:
    st.session_state.stop_event.set()
    st.warning("⏹ Arrêt demandé — fin du crawl en cours...")

if btn_start:
    if not url_input.startswith("http"):
        st.error("❌ L'URL doit commencer par http:// ou https://")
        st.stop()

    st.session_state.audit_done   = False
    st.session_state.audit_result = None
    stop_ev = threading.Event()
    st.session_state.stop_event = stop_ev
    st.session_state.config     = config

    st.markdown('<p class="section-title">🕷️ Crawl en cours</p>', unsafe_allow_html=True)

    col_prog, col_stat = st.columns([3, 1])
    progress_bar  = col_prog.progress(0)
    stat_text     = col_stat.empty()
    log_container = st.empty()
    log_lines     = []

    def on_progress(n, total):
        progress_bar.progress(min(n / total, 1.0))
        stat_text.metric("Pages", f"{n:,}/{total:,}")

    def on_log(msg):
        log_lines.append(msg)
        # Garder les 60 dernières lignes
        display = log_lines[-60:]
        log_container.markdown(
            f'<div class="crawl-log">{"<br>".join(display)}</div>',
            unsafe_allow_html=True
        )

    with st.spinner("Crawl en cours..."):
        df_pages = crawl_site_threaded(
            config,
            progress_cb=on_progress,
            log_cb=on_log,
            stop_event_ext=stop_ev,
        )

    st.session_state.df_pages = df_pages
    progress_bar.progress(1.0)

    if df_pages.empty:
        st.error("❌ Aucune page crawlée. Vérifiez l'URL et votre connexion.")
        st.stop()

    on_log(f"\n📊 Analyse des données...")
    audit = analyze_crawl(df_pages, config)
    st.session_state.audit_result = audit
    st.session_state.audit_done   = True
    st.success(f"✅ Audit terminé — {len(df_pages):,} pages analysées · Score : {audit['score_global']}/100")
    st.rerun()


# ─── AFFICHAGE DES RÉSULTATS ─────────────────────────────────────
if st.session_state.audit_done and st.session_state.audit_result:
    audit  = st.session_state.audit_result
    config = st.session_state.config
    df     = st.session_state.df_pages

    tabs = st.tabs([
        "📊 Résumé",
        "🔧 Technique",
        "📝 On-Page",
        "⚠️ Issues",
        "📋 Données brutes",
        "📥 Exports",
    ])

    # ═══ TAB 1 — RÉSUMÉ ════════════════════════════════════════
    with tabs[0]:
        st.markdown('<p class="section-title">Score global & KPIs</p>', unsafe_allow_html=True)

        col_gauge, col_kpis = st.columns([1, 2])
        with col_gauge:
            st.pyplot(make_chart_gauge(audit["score_global"]), use_container_width=True)

        with col_kpis:
            sc = audit["score_global"]
            verdict = "🟢 Excellent" if sc>=80 else "🟡 Satisfaisant" if sc>=60 else "🟠 À améliorer" if sc>=40 else "🔴 Critique"
            st.markdown(f"### {verdict}")
            st.markdown(f"**{audit['total_pages']:,}** pages analysées sur `{config['url']}`")
            st.markdown("")

            c1,c2,c3,c4 = st.columns(4)
            c1.metric("Pages OK",       f"{audit['pages_200']:,}")
            c2.metric("Erreurs",        f"{audit['pages_errors']:,}")
            c3.metric("Redirections",   f"{audit['pages_redirect']:,}")
            c4.metric("Tps moyen",      f"{audit['avg_load_time']}s")
            c5,c6,c7,c8 = st.columns(4)
            c5.metric("Noindex",        f"{audit['noindex_count']:,}")
            c6.metric("Sans canonical", f"{audit['no_canonical']:,}")
            c7.metric("Schema.org",     f"{audit['pages_schema']:,}")
            c8.metric("Pages HTTPS",    f"{audit['pages_https']:,}")

        st.markdown("---")
        st.markdown('<p class="section-title">Scores par axe SEO</p>', unsafe_allow_html=True)
        st.pyplot(make_chart_scores(audit["scores"]), use_container_width=True)

        st.markdown("---")
        st.markdown('<p class="section-title">Recommandations prioritaires</p>', unsafe_allow_html=True)

        actions = sorted([
            ("Titles",   audit["scores"]["Titles"],   "Corriger les titles manquants ou hors limites",    "Élevé"),
            ("Meta",     audit["scores"]["Meta desc"], "Rédiger des meta descriptions uniques par page",   "Élevé"),
            ("H1",       audit["scores"]["H1"],        "S'assurer d'exactement 1 H1 par page",             "Élevé"),
            ("Images",   audit["scores"]["Images alt"],"Ajouter l'attribut alt sur toutes les images",     "Moyen"),
            ("Vitesse",  audit["scores"]["Vitesse"],   "Optimiser les temps de chargement",                "Élevé"),
        ], key=lambda x:x[1])

        for i,(axe,score,action,impact) in enumerate(actions,1):
            color = "🔴" if score<40 else "🟠" if score<70 else "🟢"
            st.markdown(f"**P{i}** {color} **{axe}** (score {score}/100) — {action}")

    # ═══ TAB 2 — TECHNIQUE ═════════════════════════════════════
    with tabs[1]:
        col_pie, col_speed = st.columns(2)
        with col_pie:
            st.markdown('<p class="section-title">Codes HTTP</p>', unsafe_allow_html=True)
            st.pyplot(make_chart_http(audit["status_counts"]), use_container_width=True)
            for code, count in sorted(audit["status_counts"].items()):
                badge_cls = ("badge-200" if code==200 else
                             "badge-3xx" if code in (301,302) else
                             "badge-4xx" if code in (404,500) else "badge-0")
                st.markdown(
                    f'<span class="status-badge {badge_cls}">{code}</span> '
                    f'— **{count:,}** pages',
                    unsafe_allow_html=True
                )

        with col_speed:
            st.markdown('<p class="section-title">Temps de chargement</p>', unsafe_allow_html=True)
            if not audit["df_ok"].empty:
                st.pyplot(make_chart_speed(audit["df_ok"], config["load_time_warn"]),
                          use_container_width=True)
                st.metric("Temps moyen", f"{audit['avg_load_time']}s")
                n_slow = len(audit["issues"]["speed"])
                if n_slow:
                    st.warning(f"⚠️ {n_slow} page(s) dépassent {config['load_time_warn']}s")

        st.markdown("---")
        st.markdown('<p class="section-title">Distribution profondeur des pages</p>', unsafe_allow_html=True)
        if "depth" in df.columns:
            depth_df = df[df["status"]==200].groupby("depth").size().reset_index(name="pages")
            if not depth_df.empty:
                fig_d, ax_d = plt.subplots(figsize=(8,3))
                ax_d.bar(depth_df["depth"], depth_df["pages"],
                         color=COLORS_MAP["info"], alpha=0.8, edgecolor="white")
                ax_d.set_xlabel("Profondeur (clics depuis l'accueil)")
                ax_d.set_ylabel("Nb de pages")
                ax_d.spines[["top","right"]].set_visible(False)
                ax_d.grid(axis="y",alpha=0.3)
                fig_d.tight_layout()
                st.pyplot(fig_d, use_container_width=True)
                plt.close(fig_d)

    # ═══ TAB 3 — ON-PAGE ═══════════════════════════════════════
    with tabs[2]:
        col_t, col_m = st.columns(2)

        with col_t:
            st.markdown('<p class="section-title">Titles</p>', unsafe_allow_html=True)
            sc_t = audit["scores"]["Titles"]
            color_t = "normal" if sc_t>=70 else "inverse"
            st.metric("Score", f"{sc_t}/100", delta=f"{sc_t-70} vs seuil", delta_color=color_t)
            st.metric("Longueur moyenne", f"{audit['avg_title_len']} car.")
            if not audit["df_ok"].empty:
                st.pyplot(make_chart_titles(audit["df_ok"],
                          config["title_min"], config["title_max"]),
                          use_container_width=True)

        with col_m:
            st.markdown('<p class="section-title">Meta descriptions</p>', unsafe_allow_html=True)
            sc_m = audit["scores"]["Meta desc"]
            color_m = "normal" if sc_m>=70 else "inverse"
            st.metric("Score", f"{sc_m}/100", delta=f"{sc_m-70} vs seuil", delta_color=color_m)
            st.metric("Longueur moyenne", f"{audit['avg_meta_len']} car.")
            absent_meta = len(audit["issues"]["meta"])
            st.info(f"📊 {absent_meta} problème(s) de meta description détecté(s)")

        st.markdown("---")
        col_h1, col_img = st.columns(2)

        with col_h1:
            st.markdown('<p class="section-title">Balises H1</p>', unsafe_allow_html=True)
            sc_h = audit["scores"]["H1"]
            st.metric("Score", f"{sc_h}/100")
            no_h1    = len(audit["df_ok"][audit["df_ok"]["h1_count"]==0]) if not audit["df_ok"].empty else 0
            multi_h1 = len(audit["df_ok"][audit["df_ok"]["h1_count"]>1])  if not audit["df_ok"].empty else 0
            st.warning(f"Sans H1 : {no_h1} | Multiples H1 : {multi_h1}")

        with col_img:
            st.markdown('<p class="section-title">Images sans alt</p>', unsafe_allow_html=True)
            sc_i = audit["scores"]["Images alt"]
            st.metric("Score", f"{sc_i}/100")
            total_imgs_missing = int(audit["df_ok"]["imgs_no_alt"].sum()) if not audit["df_ok"].empty else 0
            st.warning(f"Total images sans alt : {total_imgs_missing:,}")

    # ═══ TAB 4 — ISSUES ════════════════════════════════════════
    with tabs[3]:
        issue_tabs = st.tabs(["Title","Meta","H1","Images","Vitesse"])

        for tab_el, key, label in zip(
            issue_tabs,
            ["title","meta","h1","images","speed"],
            ["Title","Meta","H1","Images","Vitesse"]
        ):
            with tab_el:
                df_iss = audit["issues"].get(key, pd.DataFrame())
                if df_iss is None or df_iss.empty:
                    st.success(f"✅ Aucun problème sur l'axe **{label}**")
                else:
                    n_crit = len(df_iss[df_iss["sévérité"]=="critique"]) if "sévérité" in df_iss.columns else 0
                    n_warn = len(df_iss) - n_crit
                    c1,c2,c3 = st.columns(3)
                    c1.metric("Total problèmes", len(df_iss))
                    c2.metric("Critiques 🔴",    n_crit)
                    c3.metric("Warnings 🟠",     n_warn)
                    st.dataframe(df_iss, use_container_width=True, height=400)
                    csv = df_iss.to_csv(index=False).encode("utf-8")
                    st.download_button(
                        f"📥 Télécharger issues {label} (.csv)",
                        data=csv,
                        file_name=f"issues_{key}.csv",
                        mime="text/csv",
                    )

    # ═══ TAB 5 — DONNÉES BRUTES ════════════════════════════════
    with tabs[4]:
        st.markdown('<p class="section-title">Données brutes du crawl</p>', unsafe_allow_html=True)

        col_f1, col_f2, col_f3 = st.columns(3)
        filter_status  = col_f1.multiselect("Filtrer par status",
                                             options=sorted(df["status"].unique()),
                                             default=[])
        filter_noindex = col_f2.checkbox("Noindex uniquement")
        filter_no_h1   = col_f3.checkbox("Sans H1 uniquement")

        df_display = df.copy()
        if filter_status:
            df_display = df_display[df_display["status"].isin(filter_status)]
        if filter_noindex:
            df_display = df_display[df_display["noindex"]==True]
        if filter_no_h1:
            df_display = df_display[df_display["h1_count"]==0]

        cols_show = ["url","status","load_time","title_len","meta_desc_len",
                     "h1_count","canonical","noindex","imgs_no_alt","has_schema"]
        cols_exist = [c for c in cols_show if c in df_display.columns]

        st.markdown(f"**{len(df_display):,}** pages affichées")
        st.dataframe(df_display[cols_exist], use_container_width=True, height=500)

        csv_all = df.to_csv(index=False).encode("utf-8")
        st.download_button(
            "📥 Télécharger CSV complet",
            data=csv_all,
            file_name="crawl_complet.csv",
            mime="text/csv",
            use_container_width=True,
        )

    # ═══ TAB 6 — EXPORTS ═══════════════════════════════════════
    with tabs[5]:
        st.markdown('<p class="section-title">Générer les rapports</p>', unsafe_allow_html=True)

        st.markdown("""
        <div class="tip-box">
        Les rapports sont générés à la demande. Cliquez sur un bouton ci-dessous pour
        créer et télécharger immédiatement le fichier.
        </div>
        """, unsafe_allow_html=True)

        col_pdf, col_pptx = st.columns(2)

        with col_pdf:
            st.markdown("#### 📄 Rapport PDF")
            st.markdown("Rapport complet — couverture, KPIs, graphiques, tableaux d'issues, plan d'action.")
            if st.button("Générer le PDF", use_container_width=True, type="primary"):
                with st.spinner("Génération du PDF..."):
                    try:
                        pdf_bytes = generate_pdf_report(audit, config)
                        st.download_button(
                            "⬇️ Télécharger le PDF",
                            data=pdf_bytes,
                            file_name=f"audit_seo_{urlparse(config['url']).netloc}_{datetime.now().strftime('%Y%m%d')}.pdf",
                            mime="application/pdf",
                            use_container_width=True,
                        )
                        st.success("✅ PDF prêt !")
                    except Exception as e:
                        st.error(f"Erreur PDF : {e}")

        with col_pptx:
            st.markdown("#### 📊 Google Slides (PPTX)")
            st.markdown("Présentation client — couverture, KPIs, scores, issues par axe, plan d'action.")
            if st.button("Générer les Slides", use_container_width=True, type="primary"):
                with st.spinner("Génération des slides..."):
                    try:
                        pptx_bytes = generate_pptx_report(audit, config)
                        st.download_button(
                            "⬇️ Télécharger le PPTX",
                            data=pptx_bytes,
                            file_name=f"audit_seo_{urlparse(config['url']).netloc}_{datetime.now().strftime('%Y%m%d')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                        st.success("✅ Slides prêtes ! Importez dans Google Slides via Fichier → Importer.")
                    except Exception as e:
                        st.error(f"Erreur PPTX : {e}")

        st.markdown("---")
        st.markdown("#### 📦 Export CSV")
        csv_export = df.to_csv(index=False).encode("utf-8")
        st.download_button(
            "⬇️ Télécharger les données brutes (.csv)",
            data=csv_export,
            file_name=f"crawl_{urlparse(config['url']).netloc}.csv",
            mime="text/csv",
            use_container_width=True,
        )

else:
    # Écran d'accueil avant le premier audit
    st.markdown("""
    <div style="text-align:center;padding:3rem 1rem;color:#888">
      <div style="font-size:4rem;margin-bottom:1rem">🔍</div>
      <h3 style="color:#455a64">Prêt à auditer</h3>
      <p>Renseignez l'URL du site dans la barre latérale et cliquez sur <strong>Lancer l'audit</strong>.</p>
    </div>
    """, unsafe_allow_html=True)

    with st.expander("ℹ️ Comment lancer l'application"):
        st.code("""
# 1. Installer les dépendances
pip install streamlit requests beautifulsoup4 pandas \\
            matplotlib reportlab python-pptx lxml

# 2. Lancer l'application
streamlit run seo_audit_app.py

# 3. Ouvrir dans le navigateur
# → http://localhost:8501
        """, language="bash")

    with st.expander("☁️ Déployer sur Streamlit Cloud (gratuit)"):
        st.markdown("""
1. Créez un dépôt GitHub avec `seo_audit_app.py` et un fichier `requirements.txt`
2. Dans `requirements.txt` :
```
streamlit
requests
beautifulsoup4
pandas
matplotlib
reportlab
python-pptx
lxml
numpy
```
3. Connectez-vous sur [share.streamlit.io](https://share.streamlit.io)
4. Déployez en 2 clics — URL publique partageable avec vos clients
        """)
